from flask import Flask, render_template, request, jsonify, send_from_directory
import requests
from pptx import Presentation
import os

app = Flask(__name__)

# Function to send text to the AIPipe API
def generate_ppt_content(text, api_key):
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"  # Adding content type as application/json
    }

    # AIPipe OpenRouter API URL
    api_url = "https://aipipe.org/openrouter/v1/chat/completions"

    data = {
        "model": "google/gemini-2.0-flash-lite-001",  # Example model, change to required one
        "messages": [{"role": "user", "content": text}]
    }

    # Send the request to the AIPipe API
    response = requests.post(api_url, headers=headers, json=data)

    if response.status_code == 200:
        return response.json()  # Return API response if successful
    else:
        return {"error": response.text}  # Return error message if something goes wrong

# Function to create slides
def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.shapes.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Function to generate PowerPoint from API response
def generate_ppt_from_api_response(api_response):
    prs = Presentation()

    content = api_response.get('choices', [{}])[0].get('message', {}).get('content', '')

    if content:
        # Split the content into sections based on double newlines
        sections = content.split('\n\n')

        # Generate a slide for each section
        for idx, section in enumerate(sections):
            create_slide(prs, f"Slide {idx + 1}", section)

    # Save the PowerPoint presentation
    ppt_filename = 'generated_presentation.pptx'
    prs.save(ppt_filename)
    return ppt_filename

@app.route('/', methods=['GET', 'POST'])
def home():
    if request.method == 'POST':
        text_input = request.form['text']  # User text input
        api_key = request.form['api_key']  # User API key from the form

        # Call the AIPipe API to get the response
        result = generate_ppt_content(text_input, api_key)

        # Generate PowerPoint from the API response
        ppt_file = generate_ppt_from_api_response(result)

        # Return the filename to the template
        return render_template('index.html', file=ppt_file)  # Pass file to the template

    return render_template('index.html')

# Route to download the generated PPT
@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(os.getcwd(), filename)

if __name__ == '__main__':
    app.run(debug=True)
